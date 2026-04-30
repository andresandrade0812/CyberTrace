# ruff: noqa: E501, B905

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

ROOT = Path(__file__).resolve().parents[1]
DELIVERABLES_DIR = ROOT / "deliverables"
SUBMISSION_DIR = ROOT / "submission" / "CyberTrace_Submission"
ZIP_PATH = ROOT / "submission" / "CyberTrace_Submission.zip"

CLUSTERS_PATH = ROOT / "reports" / "real_clusters.csv"
FIGURE_PATH = ROOT / "reports" / "figures" / "nsl_kdd_clusters.png"
GROUP_MEMBERS = [
    (
        "Victor Andrade",
        [
            "Set up the Python project structure and CLI workflow.",
            "Implemented dataset preparation, clustering, tests, and packaging scripts.",
            "Integrated the real dataset analysis outputs, figures, and submission materials.",
        ],
    ),
    (
        "Ben Miller",
        [
            "Helped shape the project direction around cybersecurity investigation.",
            "Contributed to the problem framing, motivation, and interpretation goals.",
            "Supported the final narrative about findings, implications, and future work.",
        ],
    ),
]


def main() -> None:
    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    SUBMISSION_DIR.parent.mkdir(parents=True, exist_ok=True)

    clusters = pd.read_csv(CLUSTERS_PATH)
    summary = build_summary(clusters)

    report_markdown = render_report_markdown(summary)
    (DELIVERABLES_DIR / "final_report.md").write_text(report_markdown, encoding="utf-8")
    build_report_docx(summary, DELIVERABLES_DIR / "final_report.docx")
    build_presentation(summary, DELIVERABLES_DIR / "final_presentation.pptx")

    assemble_submission_package()
    make_zip_archive()

    print(f"Wrote {DELIVERABLES_DIR / 'final_report.md'}")
    print(f"Wrote {DELIVERABLES_DIR / 'final_report.docx'}")
    print(f"Wrote {DELIVERABLES_DIR / 'final_presentation.pptx'}")
    print(f"Wrote {ZIP_PATH}")


def build_summary(clusters: pd.DataFrame) -> dict[str, object]:
    label_mix = pd.crosstab(clusters["cluster"], clusters["label"])
    attack_mix = pd.crosstab(clusters["cluster"], clusters["attack_type"])

    cluster_rows: list[dict[str, object]] = []
    for cluster_id, group in clusters.groupby("cluster"):
        label_counts = group["label"].value_counts()
        dominant_label = label_counts.idxmax()
        dominant_share = label_counts.max() / len(group)
        top_attacks = group["attack_type"].value_counts().head(3)
        cluster_rows.append(
            {
                "cluster": int(cluster_id),
                "size": int(len(group)),
                "benign": int(label_counts.get("benign", 0)),
                "malicious": int(label_counts.get("malicious", 0)),
                "dominant_label": dominant_label,
                "dominant_share": dominant_share,
                "top_attacks": list(zip(top_attacks.index.tolist(), top_attacks.tolist())),
            }
        )

    attack_counts = clusters["attack_type"].value_counts()

    return {
        "record_count": int(len(clusters)),
        "benign_count": int((clusters["label"] == "benign").sum()),
        "malicious_count": int((clusters["label"] == "malicious").sum()),
        "cluster_count": int(clusters["cluster"].nunique()),
        "cluster_rows": cluster_rows,
        "label_mix": label_mix,
        "attack_mix": attack_mix,
        "top_attack_counts": attack_counts.head(8).to_dict(),
    }


def render_report_markdown(summary: dict[str, object]) -> str:
    cluster_rows = summary["cluster_rows"]
    lines = [
        "# CyberTrace Final Report",
        "",
        "**Course Project:** CS544",
        "",
        "**Group Members:** Victor Andrade and Ben Miller",
        "",
        "## Abstract",
        "",
        "CyberTrace explores how clustering can support cybersecurity analysis when suspicious",
        "traffic or artifacts do not match a known signature exactly. The project emphasizes",
        "investigative workflow rather than only classifier accuracy. In the final implementation,",
        "Python is used to prepare a safe public intrusion-detection benchmark, convert it into a",
        "numeric feature table, cluster the samples, and generate interpretable outputs for review.",
        "The real dataset milestone uses a small NSL-KDD training sample with 1,011 records and",
        "produces three clusters that separate mostly benign traffic from two malicious-heavy groups.",
        "These results suggest that clustering can help analysts prioritize review, compare related",
        "behaviors, and distinguish clearly malicious traffic from normal-like activity.",
        "",
        "## 1. Problem Statement",
        "",
        "The problem addressed by CyberTrace is how to compare suspicious network-related artifacts",
        "without depending only on exact signature matching. Traditional rule-based and signature-based",
        "defenses remain valuable, but they can miss new, modified, or behaviorally subtle threats.",
        "Analysts still need a way to examine observable characteristics and organize suspicious activity",
        "into groups that support investigation.",
        "",
        "More specifically, this project asks whether safe network traffic records can be transformed into",
        "a useful feature representation and then clustered into groups that reveal meaningful differences",
        "between benign and malicious behavior. The goal is not to claim perfect detection, but to build",
        "a reproducible analysis workflow that helps answer why certain samples look suspicious and how",
        "they relate to one another.",
        "",
        "A useful final project in this space should therefore do three things well. First, it should",
        "represent the problem in a way that is close to actual analyst reasoning. Second, it should be",
        "implemented in a reproducible way rather than as a fragile one-time experiment. Third, it should",
        "produce findings that can be discussed in plain cybersecurity language, not only in model-centric",
        "terms. CyberTrace was designed around those three needs.",
        "",
        "## 2. Motivation",
        "",
        "This problem is important because modern cybersecurity work involves more than checking whether",
        "something matches a known bad indicator. Analysts often face large sets of events, connections,",
        "or artifacts that need quick prioritization. If clustering can place similar behaviors together,",
        "then it can reduce analyst effort, reveal repeated suspicious patterns, and provide a clearer",
        "starting point for incident response or deeper forensic review.",
        "",
        "The motivation is also practical for a course project. A safe benchmark dataset can still capture",
        "many of the behavioral questions that matter in cybersecurity: error rates, destination patterns,",
        "service usage, and connection counts. By focusing on those interpretable signals, the project",
        "stays aligned with real analyst thinking instead of turning into a purely abstract machine",
        "learning exercise.",
        "",
        "This motivation also connects directly to incident response. In a live environment, defenders",
        "often need to sort through events that are not immediately conclusive. A clustering-based view",
        "does not replace expert judgment, but it can create structure around that judgment by showing",
        "which samples appear isolated, which appear repeated, and which share behavior with clearly",
        "malicious records. That makes the exploration process more systematic.",
        "",
        "## 3. Dataset and Safety",
        "",
        "The final project uses a small public NSL-KDD training CSV sample stored in",
        "`data/raw/nsl_kdd_small_training.csv`. NSL-KDD is a safe, labeled intrusion-detection",
        "benchmark derived from KDD-style connection records. It is appropriate for this course",
        "project because it avoids handling live malware or sensitive production traffic while still",
        "providing realistic benign and attack labels for interpretation.",
        "",
        f"- Total records: {summary['record_count']}",
        f"- Benign records: {summary['benign_count']}",
        f"- Malicious records: {summary['malicious_count']}",
        "- Source mirror: Jehuty4949/NSL_KDD GitHub repository",
        "- Benchmark reference: Canadian Institute for Cybersecurity NSL-KDD page",
        "",
        "Using NSL-KDD lets the project focus on exploratory cybersecurity analysis without introducing",
        "the risk of storing live malware, sensitive enterprise packet captures, or private credentials",
        "inside the repository. That safety constraint matters because the project should be reproducible,",
        "shareable, and suitable for course submission.",
        "",
        "The dataset is also useful because it includes a mixture of normal traffic and several attack",
        "categories such as neptune, portsweep, ipsweep, and satan. Even though NSL-KDD is older, it still",
        "provides a reasonable environment for discussing the core question of this project: whether groups",
        "of records can be organized into meaningful behavioral clusters that reflect benign or malicious",
        "activity.",
        "",
        "## 4. Contributions",
        "",
        "The main contributions of this work are:",
        "",
        "1. A Python repository structure designed around cybersecurity investigation workflows.",
        "2. A command-line data preparation step for converting raw NSL-KDD rows into a CyberTrace",
        "   feature table with metadata and numeric model inputs.",
        "3. A clustering pipeline that standardizes features, applies KMeans, and writes cluster",
        "   assignments for further analysis.",
        "4. Notebooks, figures, tests, and packaging scripts that make the project reproducible.",
        "5. A final interpretation of how the discovered clusters relate to benign and malicious behavior.",
        "",
        "An important contribution is that the project does not stop at code execution. It also packages",
        "the work into a report, presentation, and submission archive so the analysis can be reviewed by",
        "someone who did not watch the implementation happen. That makes the work more complete and more",
        "appropriate for a final academic submission.",
        "",
        "## 5. Methodology",
        "",
        "The project workflow has four main stages:",
        "",
        "1. Store raw safe datasets in `data/raw/`.",
        "2. Prepare the dataset into a CyberTrace feature table with the `prepare-nsl-kdd` command.",
        "3. Standardize numeric features and cluster them with KMeans.",
        "4. Interpret cluster membership by comparing labels and attack types after clustering.",
        "",
        "The NSL-KDD preparer converts raw rows into a table containing metadata columns",
        "(`sample_id`, `label`, `attack_type`, `source`, and `difficulty`) plus numeric features and",
        "one-hot encoded categorical protocol, service, and flag fields. Clustering is then applied",
        "only to the numeric feature matrix, not to the labels.",
        "",
        "This approach matters because it preserves the separation between unsupervised learning and",
        "interpretation. The labels are not used to fit the clustering model; instead, they are used",
        "afterward to understand what kinds of behaviors ended up together. That is much closer to how",
        "clustering would be used in an exploratory cybersecurity workflow.",
        "",
        "KMeans was chosen as the baseline method because it is easy to explain, easy to reproduce, and",
        "useful for a first clustering pass over standardized numeric features. The cluster count of three",
        "was selected as a straightforward baseline that could separate normal-like, mixed, and strongly",
        "malicious behavior. This is not presented as the only correct choice, but as an interpretable",
        "starting point for the exploration.",
        "",
        "The methodology also includes validation at the software level. The repository contains automated",
        "tests for clustering, feature loading, extraction, and dataset preparation. That matters because",
        "a cybersecurity analysis workflow should not depend only on manually repeated notebook steps. By",
        "testing the code paths that produce the prepared data and final outputs, the project improves its",
        "reliability and repeatability.",
        "",
        "From an investigative standpoint, the feature preparation step is especially important. Raw traffic",
        "records are rarely useful unless they can be summarized into stable technical signals. In this",
        "project, those signals include byte-oriented values, counts, error-related rates, service usage,",
        "and protocol or flag indicators. Each of these features captures one piece of behavioral evidence,",
        "and clustering works by evaluating how those pieces line up across many records at once.",
        "",
        "This feature-centric view is one reason the project remains explainable. Even if the clustering",
        "itself is unsupervised, the inputs still have a concrete meaning that analysts can discuss. That",
        "means the final write-up can talk about repeated attack behavior, normal-like connections, service",
        "patterns, and suspicious concentration, rather than presenting the clusters as unexplained black",
        "box outputs.",
        "",
        "Another methodological point is the choice to keep the workflow simple enough to inspect end to",
        "end. A more complex model or feature-engineering pipeline might produce different outputs, but it",
        "would also make the final project harder to explain. For this assignment, interpretability and",
        "clarity were treated as design constraints. That is why the project prefers a direct preparation",
        "step, standardized numeric features, and an unsupervised model that can be described cleanly in",
        "both code and prose.",
        "",
        "This also affects how the results should be read. The clusters are not presented as final truths",
        "about attack behavior. They are analytical groupings produced under a specific feature design and",
        "specific clustering choice. Their value lies in how well they support comparison and reasoning, not",
        "in the claim that they are the only valid partition of the data. That distinction is important in",
        "cybersecurity, where useful investigation often depends on creating workable structure rather than",
        "finding one perfectly correct answer.",
        "",
        "## 6. Implementation",
        "",
        "CyberTrace is implemented as a Python package under `src/cybertrace/`. The main commands are:",
        "",
        "- `cybertrace prepare-nsl-kdd data/raw/nsl_kdd_small_training.csv --output data/processed/features_nsl_kdd.csv`",
        "- `cybertrace cluster data/processed/features_nsl_kdd.csv --output reports/real_clusters.csv --clusters 3`",
        "",
        "The codebase also includes starter synthetic data, a simple network-log extractor, notebooks,",
        "tests, and figures. The repository was validated with Ruff and Pytest before packaging.",
        "",
        "The implementation also includes a reproducible submission builder that generates the final",
        "report, presentation slides, and zip archive. This makes the final project easier to inspect,",
        "rerun, and update if small report changes are needed before submission.",
        "",
        "From a software-engineering perspective, the project is organized into clear responsibilities.",
        "Dataset preparation logic lives in a dataset module, clustering logic lives in its own module,",
        "shared configuration is centralized, and the command-line interface ties those pieces together.",
        "This structure makes the code easier to extend to new datasets or additional analysis commands.",
        "",
        "The use of notebooks is intentionally limited to exploration and communication. The core logic",
        "needed to reproduce the project results lives in the package and CLI commands, which means the",
        "project can still be rerun from the terminal without depending on interactive notebook state.",
        "",
        "That design choice has practical benefits for a final project. It means someone reviewing the",
        "submission can run the code paths directly, inspect the generated artifacts, and understand the",
        "main flow without reconstructing hidden notebook state. It also makes the project more maintainable",
        "because new datasets or additional reporting layers can be added incrementally as separate modules.",
        "",
        "## 7. Findings and Results",
        "",
        "The real dataset clustering produced three groups:",
        "",
    ]

    for row in cluster_rows:
        lines.append(
            f"- Cluster {row['cluster']}: {row['size']} samples "
            f"({row['benign']} benign, {row['malicious']} malicious, "
            f"{row['dominant_share']:.1%} {row['dominant_label']})"
        )

    lines.extend(
        [
            "",
            "The cluster label mix is:",
            "",
            f"- Cluster 0: {summary['cluster_rows'][0]['benign']} benign and {summary['cluster_rows'][0]['malicious']} malicious",
            f"- Cluster 1: {summary['cluster_rows'][1]['benign']} benign and {summary['cluster_rows'][1]['malicious']} malicious",
            f"- Cluster 2: {summary['cluster_rows'][2]['benign']} benign and {summary['cluster_rows'][2]['malicious']} malicious",
            "",
        "The most common attack types in the sample are:",
        "",
        ]
    )

    for attack, count in summary["top_attack_counts"].items():
        lines.append(f"- {attack}: {count}")

    lines.extend(
        [
            "",
            "### Cluster Interpretation",
            "",
            "Cluster 0 is primarily benign and functions as the most normal-like group in the analysis.",
            "It still contains a smaller malicious subset, which reflects the fact that unsupervised",
            "clustering groups by shared behavior rather than by ground-truth labels alone.",
            "",
            "Cluster 1 is malicious-heavy and includes a mixture of neptune, portsweep, satan, and a",
            "smaller number of benign records. This suggests a mixed suspicious-behavior cluster where",
            "attack traffic shares count-based or rate-based characteristics with some normal connections.",
            "",
        "Cluster 2 is entirely malicious in this sample and is dominated by neptune records. This",
        "is the clearest malicious cluster in the baseline run and shows that clustering can isolate",
        "a strongly attack-oriented region of the feature space.",
        "",
        "These results reveal an important pattern. The baseline run did not create one single malicious",
        "cluster containing every attack record. Instead, it produced at least two malicious-heavy groups:",
        "one mixed cluster and one extremely concentrated malicious cluster. That suggests the malicious",
        "traffic in the sample is not behaviorally uniform. Some attacks share characteristics with each",
        "other closely enough to form a concentrated cluster, while others overlap more with a wider range",
        "of traffic patterns.",
        "",
        "That distinction matters for cybersecurity interpretation. A strongly concentrated malicious",
        "cluster may correspond to repeated or stereotyped attack behavior, which is often easier to flag",
        "and isolate. A mixed malicious-heavy cluster is different: it reflects a region of the feature",
        "space where suspicious traffic may blend with more ordinary behavior, making investigation more",
        "challenging and more dependent on contextual review.",
        "",
        "The attack-type distribution helps reinforce this point. The sample contains a substantial number",
        "of neptune records, and those records dominate the most malicious cluster. By contrast, portsweep,",
        "ipsweep, satan, and smaller attack categories appear in more mixed groupings. That suggests some",
        "attack behaviors are much more behaviorally concentrated than others, which is exactly the kind of",
        "pattern a clustering-based exploratory workflow should be able to reveal.",
        "",
        "The benign records also tell an important story. Most benign traffic is concentrated in Cluster 0,",
        "which gives the analysis a clear normal-like reference point. However, a smaller number of benign",
        "records appear in a malicious-heavy cluster. That overlap is a reminder that not every suspicious",
        "looking pattern is malicious and that not every benign record is behaviorally simple. In practice,",
        "that is why analysts need both grouping methods and contextual judgment.",
        "",
        "Taken together, the results suggest that clustering is most valuable here as an organizational aid.",
        "It helps reveal where the data is strongly separable and where it is not. Strongly separable regions",
        "can help identify concentrated suspicious behavior quickly, while mixed regions indicate where more",
        "careful investigation or additional evidence is required. In other words, the output is useful not",
        "only for what it separates, but also for what it leaves ambiguous.",
        "",
        "An important takeaway is that clustering did not create a perfect benign-versus-malicious split,",
        "and that is actually useful to discuss. Unsupervised methods group by similarity, so mixed",
        "clusters highlight where suspicious and benign records share overlapping characteristics. That",
        "kind of ambiguity is realistic in cybersecurity analysis and is often where investigators spend",
        "the most time.",
        "",
        "## 8. Implications",
            "",
            "The results support the project goal: clustering can help analysts organize traffic into useful",
            "groups and interpret those groups in cybersecurity terms. The workflow does not replace",
            "signature-based or supervised detection, but it creates a structured investigative view of the",
            "data. In practice, this can help analysts prioritize clusters for deeper review, identify",
            "families of suspicious behavior, and compare what separates benign from malicious traffic",
            "patterns.",
            "",
        "Another implication is educational: a relatively compact Python project can still model the full",
        "arc of a cybersecurity data analysis task. The repo now includes preparation, experimentation,",
        "interpretation, reporting, and packaging, which makes it a strong demonstration project beyond",
        "just a one-off notebook.",
        "",
        "There is also an implication for how students and practitioners think about unsupervised methods.",
        "Clustering is sometimes treated as a vague exploratory add-on, but in this project it becomes a",
        "practical organizational tool. The value is not only in the cluster assignments themselves, but",
        "in the discipline of preparing interpretable features, comparing group composition, and asking",
        "what characteristics make one cluster look more suspicious than another.",
        "",
        "The project also suggests a broader lesson about cybersecurity analytics: interpretation is often",
        "the bridge between data science and operational usefulness. A model output by itself is not enough.",
        "What makes the output valuable is the ability to explain why a cluster appears suspicious, what",
        "technical signals are likely contributing to that grouping, and how an analyst might act on the",
        "finding. CyberTrace intentionally emphasizes that interpretive step.",
        "",
        "Because of that emphasis, the project contributes not only a result but also a workflow pattern.",
        "The pattern is: prepare safe data, cluster behaviorally similar records, compare discovered groups",
        "to known labels after the fact, and then translate those groups into plain-language cybersecurity",
        "insights. That pattern is portable and can be reused on other safe datasets or more advanced future",
        "versions of the project.",
        "",
        "A final implication is that cybersecurity education benefits from projects that combine technical",
        "implementation with interpretive responsibility. It is one thing to produce a cluster assignment",
        "file. It is another to explain what the clusters suggest, what they fail to capture, and why an",
        "analyst should trust or question the result. CyberTrace is strongest where it treats those tasks as",
        "part of the same project rather than as separate concerns.",
        "",
        "## 9. Group Members and Contributions",
            "",
            "The project was completed by Victor Andrade and Ben Miller.",
            "",
            "### Victor Andrade",
            "",
            "- Implemented the Python repository structure and CLI commands.",
            "- Added dataset preparation, clustering, tests, notebooks, and deliverable packaging.",
            "- Integrated the real dataset outputs and final submission materials.",
            "",
            "### Ben Miller",
            "",
            "- Helped shape the cybersecurity framing of the project.",
            "- Contributed to the problem definition, motivation, and interpretation goals.",
        "- Supported the final explanation of findings, implications, and future directions.",
        "",
        "These contribution descriptions are meant to make the division of work visible in the final",
        "submission rather than leaving the project to appear as a single undifferentiated artifact. If",
        "the group wants to revise the exact wording, this section can be edited easily before submission.",
        "",
        "## 10. Limitations",
            "",
            "- NSL-KDD is an older benchmark and may not fully represent modern enterprise traffic.",
            "- The current clustering uses KMeans with a fixed cluster count of three for the baseline run.",
            "- This project focuses on network-connection style benchmark features rather than live packet",
            "  parsing or current malware families.",
        "- Interpretation is aided by labels after clustering, so future work should add stronger",
        "  feature-importance analysis and expanded datasets.",
        "",
        "Another limitation is that the current report emphasizes one baseline clustering result rather",
        "than a larger sweep across many parameter settings. That is acceptable for a course project, but",
        "a deeper study would compare alternative values of k, different random seeds, and potentially",
        "different distance-based or density-based clustering methods.",
        "",
        "## 11. Future Work",
            "",
            "- Evaluate additional safe public datasets such as CIC-style flow data.",
            "- Compare KMeans with hierarchical clustering or DBSCAN.",
            "- Add packet-capture or Zeek/Suricata preparation commands for richer lab artifacts.",
            "- Expand reporting with more visualizations and cluster-specific feature summaries.",
            "",
            "Future work should also focus on more modern datasets and a stronger comparison of clustering",
            "quality. For example, silhouette scores, per-feature summary tables, and comparisons across",
            "multiple values of k would help show how stable the discovered groups are. Another strong next",
        "step would be mapping the cluster outputs to concrete analyst actions, such as escalation,",
        "triage, or additional artifact collection.",
        "",
        "Future work could also connect the project more directly to operational workflows. For example,",
        "cluster outputs might be used to generate short analyst notes, candidate incident categories, or",
        "recommendations for which samples deserve manual review first. That would move the project from",
        "a strong exploratory analysis into a stronger analyst-support prototype.",
        "",
        "## 12. Reproducibility and Practical Value",
        "",
        "One of the strengths of the final project is that the entire workflow can be rerun from the",
        "repository itself. The raw dataset, preparation step, clustering command, notebooks, figures,",
        "tests, report source, slide deck, and submission builder are all included together. That means",
        "the project is not just a set of final screenshots or manually written conclusions; it is a",
        "working pipeline that can be inspected and executed again.",
        "",
        "This reproducibility has practical value for cybersecurity work. Analysts and teams often need to",
        "review not only results, but also how those results were produced. A clear chain from raw input",
        "to prepared features to cluster assignments to interpretation makes it easier to trust the output,",
        "spot weaknesses, and extend the workflow to new datasets. Even in a classroom context, that kind",
        "of transparency is worth emphasizing because it reflects good security engineering practice.",
        "",
        "The project is also useful as a template. Although the current implementation centers on a safe",
        "benchmark, the same structure could be reused for lab packet captures, Zeek logs, Suricata alerts,",
        "or other safe feature-rich artifacts. In that sense, the final submission is not only a finished",
        "assignment but also a foundation for future cybersecurity analysis projects.",
        "",
        "## 13. Conclusion",
        "",
        "CyberTrace demonstrates that a Python-based clustering pipeline can support cybersecurity",
        "analysis by preparing safe datasets, grouping samples by behavior, and producing interpretable",
        "outputs. The final implementation achieves the project objective by delivering code, data,",
            "results, tests, a report, and presentation materials in a reproducible repository.",
        ]
    )
    return "\n".join(lines) + "\n"


def build_report_docx(summary: dict[str, object], output_path: Path) -> None:
    document = Document()
    style = document.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("CyberTrace Final Report")
    run.bold = True
    run.font.size = Pt(18)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Clustering and Analysis of Suspicious Network Traffic").italic = True
    members = document.add_paragraph()
    members.alignment = WD_ALIGN_PARAGRAPH.CENTER
    members.add_run("Group Members: Victor Andrade and Ben Miller")

    add_heading(document, "Abstract")
    document.add_paragraph(
        "CyberTrace is a cybersecurity analysis project that uses Python to prepare a safe public "
        "intrusion-detection dataset, cluster network behavior records, and interpret the resulting "
        "groups from an analyst perspective. Using a small NSL-KDD sample with 1,011 records, the "
        "project produced three clusters that separated mostly benign traffic from two malicious-heavy "
        "groups. The final workflow shows that clustering can support investigation beyond simple "
        "signature matching by organizing records into interpretable groups."
    )

    add_heading(document, "Problem Statement")
    document.add_paragraph(
        "The problem addressed by CyberTrace is how to compare suspicious network-related artifacts "
        "without depending only on exact signature matching. In cybersecurity practice, analysts "
        "often need to decide whether network behavior looks normal, suspicious, or malicious based "
        "on observable characteristics rather than perfect known-bad matches."
    )
    document.add_paragraph(
        "This project therefore asks whether a safe benchmark dataset can be transformed into a "
        "feature representation that supports unsupervised clustering and useful interpretation. "
        "The goal is not to replace detection systems, but to create a workflow that helps explain "
        "what samples have in common and why certain groups appear suspicious."
    )
    document.add_paragraph(
        "A useful final project in this space should represent the problem in a way that resembles "
        "real analyst reasoning, be reproducible as software, and produce results that can be "
        "explained in cybersecurity language rather than only in abstract machine learning terms. "
        "CyberTrace was designed with those priorities in mind."
    )

    add_heading(document, "Motivation")
    document.add_paragraph(
        "This problem matters because analysts routinely face large volumes of events and cannot "
        "manually investigate everything. A clustering-based view can reduce that burden by grouping "
        "similar records together, allowing analysts to focus on suspicious clusters first."
    )
    document.add_paragraph(
        "The project is also motivated by the need for explainable, analyst-friendly workflows. "
        "Rather than centering the project entirely on predictive performance, CyberTrace emphasizes "
        "cybersecurity interpretation, reproducibility, and safe handling of data."
    )
    document.add_paragraph(
        "This motivation is practical for incident response. When analysts face many events that are "
        "not immediately conclusive, they need structure. Clustering can help provide that structure "
        "by showing which records belong together, which behaviors repeat, and which groups should be "
        "reviewed first."
    )

    add_heading(document, "Dataset and Safety")
    document.add_paragraph(
        "The final project uses a small public NSL-KDD training CSV sample. This benchmark is safe "
        "for coursework because it provides labeled benign and attack records without requiring live "
        "malware samples or sensitive production traffic."
    )
    for bullet in [
        f"Total records: {summary['record_count']}",
        f"Benign records: {summary['benign_count']}",
        f"Malicious records: {summary['malicious_count']}",
        "Source mirror: Jehuty4949/NSL_KDD GitHub repository",
        "Benchmark reference: Canadian Institute for Cybersecurity NSL-KDD page",
    ]:
        document.add_paragraph(bullet, style="List Bullet")
    document.add_paragraph(
        "This choice keeps the project safe for classroom use while still supporting meaningful "
        "discussion of benign and malicious behavior patterns."
    )
    document.add_paragraph(
        "The dataset also includes several attack categories, allowing the project to discuss not only "
        "whether traffic is malicious, but whether malicious behavior itself splits into more than one "
        "kind of cluster."
    )

    add_heading(document, "Contributions")
    for bullet in [
        "Built a Python repository and command-line workflow for cybersecurity clustering analysis.",
        "Added a real dataset preparation step for NSL-KDD using reproducible code.",
        "Implemented clustering, notebook analysis, figures, tests, and submission packaging.",
        "Produced an interpreted real-data result that distinguishes normal-like and malicious-heavy clusters.",
    ]:
        document.add_paragraph(bullet, style="List Bullet")
    document.add_paragraph(
        "A further contribution is that the project packages the work into reusable deliverables. "
        "Instead of stopping at code execution, it creates a final report, slide deck, and submission "
        "zip, making the entire project easier to inspect and hand in."
    )

    add_heading(document, "Methodology")
    for step in [
        "Store safe raw datasets in data/raw/.",
        "Prepare NSL-KDD records into a numeric CyberTrace feature table.",
        "Standardize features and cluster them with KMeans.",
        "Interpret the clusters by comparing labels and attack types after clustering.",
    ]:
        document.add_paragraph(step, style="List Number")
    document.add_paragraph(
        "The preparer keeps metadata such as sample identifiers, attack type, source, and difficulty "
        "while transforming the usable model inputs into numeric columns. Categorical protocol, "
        "service, and flag fields are one-hot encoded so they can be included in the clustering step."
    )
    document.add_paragraph(
        "This separation matters because labels are used only after clustering for interpretation. "
        "That keeps the analysis unsupervised while still allowing the project to explain what kinds "
        "of behaviors each cluster contains."
    )
    document.add_paragraph(
        "KMeans was selected as the baseline clustering method because it is easy to explain, easy to "
        "reproduce, and suitable for a first pass over standardized features. The baseline run uses "
        "three clusters so the results can be interpreted in a straightforward way."
    )
    document.add_paragraph(
        "The methodology also includes software validation. Tests cover the main preparation, "
        "feature-loading, extraction, and clustering paths so the workflow is repeatable and less "
        "dependent on ad hoc notebook execution."
    )
    document.add_paragraph(
        "From an investigative standpoint, the feature preparation step is critical. Raw network-style "
        "records are rarely useful unless they can be summarized into stable technical signals. In this "
        "project, those signals include byte counts, rates, service indicators, and protocol or flag "
        "fields. Clustering then evaluates how those signals align across the full dataset."
    )
    document.add_paragraph(
        "This feature-centric view keeps the project interpretable. Even if the clustering model is "
        "unsupervised, the inputs still have concrete cybersecurity meaning, which makes the final "
        "analysis easier to explain and defend."
    )

    add_heading(document, "Implementation")
    document.add_paragraph(
        "CyberTrace is implemented as a Python package with CLI commands for dataset preparation, "
        "clustering, notebooks, figures, and tests. The real-data workflow is driven by two commands."
    )
    document.add_paragraph(
        "cybertrace prepare-nsl-kdd data/raw/nsl_kdd_small_training.csv --output "
        "data/processed/features_nsl_kdd.csv",
        style="Intense Quote",
    )
    document.add_paragraph(
        "cybertrace cluster data/processed/features_nsl_kdd.csv --output reports/real_clusters.csv "
        "--clusters 3",
        style="Intense Quote",
    )
    document.add_paragraph(
        "The repository also includes a network-log feature extractor, exploratory notebooks, "
        "generated figures, and a builder script that assembles the final report, presentation, "
        "and submission zip file."
    )
    document.add_paragraph(
        "This modular structure makes the code easier to extend. New datasets can be added through "
        "additional preparers, and alternative clustering or reporting methods can be introduced "
        "without rewriting the entire project."
    )
    document.add_paragraph(
        "That design choice also improves reviewability. A grader or collaborator can run the core "
        "commands directly and reproduce the key outputs without relying on hidden notebook state."
    )

    add_heading(document, "Findings and Results")
    table = document.add_table(rows=1, cols=5)
    header_cells = table.rows[0].cells
    header_cells[0].text = "Cluster"
    header_cells[1].text = "Size"
    header_cells[2].text = "Benign"
    header_cells[3].text = "Malicious"
    header_cells[4].text = "Dominant Label"
    for row in summary["cluster_rows"]:
        cells = table.add_row().cells
        cells[0].text = str(row["cluster"])
        cells[1].text = str(row["size"])
        cells[2].text = str(row["benign"])
        cells[3].text = str(row["malicious"])
        cells[4].text = f"{row['dominant_label']} ({row['dominant_share']:.1%})"

    document.add_paragraph()
    document.add_picture(str(FIGURE_PATH), width=Inches(6.5))
    caption = document.add_paragraph("Figure 1. PCA projection of NSL-KDD records colored by label.")
    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_paragraph(
        "Cluster 0 is primarily benign and acts as the most normal-like cluster in the baseline run. "
        "Cluster 1 is mixed but malicious-heavy, suggesting overlapping suspicious patterns. "
        "Cluster 2 is fully malicious in this sample and is dominated by neptune records."
    )
    document.add_paragraph(
        "These findings are useful because they show both the strengths and the limits of clustering. "
        "The strongest malicious pattern becomes clearly separated, but another cluster remains mixed, "
        "which reflects the ambiguity analysts often face in real investigations."
    )
    document.add_paragraph(
        "The results therefore suggest that malicious traffic in the sample is not behaviorally "
        "uniform. One group is strongly concentrated around repeated attack-like behavior, while "
        "another overlaps more with traffic that is closer to normal in some dimensions."
    )
    document.add_paragraph(
        "This is important because mixed clusters are not failures to be ignored. In many security "
        "settings, those mixed regions are exactly where analysts need more context, enrichment, and "
        "careful review."
    )
    document.add_paragraph(
        "The attack-type distribution reinforces this interpretation. Neptune records dominate the most "
        "malicious cluster, while portsweep, ipsweep, satan, and smaller attack categories appear more "
        "often in mixed groupings. That suggests some malicious behaviors are much more concentrated "
        "than others in the feature space."
    )
    document.add_paragraph(
        "The benign records are also informative. Most benign traffic appears in Cluster 0, which "
        "creates a useful normal-like reference point. A smaller number of benign records appear in a "
        "malicious-heavy cluster, showing that overlap and ambiguity remain part of the analysis."
    )

    add_heading(document, "Implications")
    document.add_paragraph(
        "The results support the larger project claim that clustering can help organize traffic for "
        "analyst review. Even without using labels to fit the model, the workflow produced one mostly "
        "benign cluster, one mixed suspicious cluster, and one clearly malicious cluster."
    )
    document.add_paragraph(
        "This matters because cybersecurity analysis often depends on triage and prioritization. "
        "A cluster-oriented view can help direct attention toward groups that deserve deeper review, "
        "additional indicators, or manual analyst validation."
    )
    document.add_paragraph(
        "The project also has educational implications. It demonstrates that a relatively compact "
        "Python repository can still cover the full workflow of a cybersecurity analytics project: "
        "data preparation, model application, interpretation, documentation, and packaging."
    )
    document.add_paragraph(
        "More broadly, the project argues that interpretation is the bridge between data science and "
        "operational usefulness. A cluster assignment by itself is not enough; the value comes from "
        "understanding why that group looks suspicious and what an analyst should do next."
    )
    document.add_paragraph(
        "For that reason, CyberTrace contributes not only a set of outputs but also a reusable workflow "
        "pattern: prepare safe data, cluster similar records, compare discovered groups with labels "
        "after the fact, and translate the result into plain-language cybersecurity insights."
    )

    add_heading(document, "Group Members")
    for member, contributions in GROUP_MEMBERS:
        document.add_paragraph(member, style="List Bullet")
        for contribution in contributions:
            paragraph = document.add_paragraph(contribution, style="List Bullet 2")
            paragraph.paragraph_format.left_indent = Inches(0.5)
    document.add_paragraph(
        "This section is included to make individual effort visible in the final submission instead "
        "of treating the entire project as a single anonymous artifact."
    )

    add_heading(document, "Limitations")
    for bullet in [
        "NSL-KDD is older than current enterprise network environments.",
        "The baseline run uses a fixed cluster count of three.",
        "The workflow currently emphasizes benchmark records over richer live packet parsing.",
        "Interpretation is supported by labels after clustering, so deeper feature analysis is still needed.",
    ]:
        document.add_paragraph(bullet, style="List Bullet")
    document.add_paragraph(
        "Another limitation is that this report emphasizes one baseline clustering result rather than "
        "a large comparison across many parameter settings. A more exhaustive study would compare "
        "multiple choices of k and alternative clustering methods."
    )
    add_heading(document, "Future Work")
    for bullet in [
        "Compare additional clustering methods such as hierarchical clustering or DBSCAN.",
        "Use newer flow-based datasets to better reflect modern traffic patterns.",
        "Add more analyst-facing outputs, such as per-cluster feature summaries and scoring.",
        "Expand preparation to other safe artifacts such as Zeek or Suricata logs.",
    ]:
        document.add_paragraph(bullet, style="List Bullet")
    document.add_paragraph(
        "These improvements would make the project stronger both as a cybersecurity investigation "
        "tool and as a more rigorous data-analysis comparison."
    )
    document.add_paragraph(
        "A strong next step would be to connect cluster outputs directly to analyst actions such as "
        "triage recommendations, escalation flags, or suggested follow-up evidence. That would push "
        "the project beyond exploration into more operational decision support."
    )

    add_heading(document, "Reproducibility and Practical Value")
    document.add_paragraph(
        "A major strength of the final project is that the entire workflow can be rerun directly from "
        "the repository. The raw dataset, preparation code, clustering command, notebooks, figures, "
        "tests, report, presentation, and submission zip are all generated from the same project."
    )
    document.add_paragraph(
        "This reproducibility matters in cybersecurity because results are more trustworthy when the "
        "path from input data to analyst-facing output is visible. A reviewer can inspect how records "
        "were prepared, how clustering was applied, and how the results were interpreted."
    )
    document.add_paragraph(
        "The project is also practically reusable. Although the final implementation focuses on "
        "NSL-KDD, the repository structure can be extended to other safe datasets such as flow logs, "
        "packet-derived features, Zeek outputs, or Suricata alerts."
    )

    add_heading(document, "Lessons Learned")
    document.add_paragraph(
        "One lesson from the project is that clean project structure matters almost as much as the "
        "clustering result itself. Once the data preparation, model execution, reporting, and tests "
        "were organized into clear modules, it became much easier to reason about the analysis and "
        "to regenerate the final outputs reliably."
    )
    document.add_paragraph(
        "A second lesson is that mixed results are still meaningful results. The presence of a mixed "
        "malicious-heavy cluster does not weaken the project; instead, it highlights the fact that "
        "real cybersecurity analysis often lives in gray areas rather than in perfect separations."
    )
    document.add_paragraph(
        "A third lesson is that packaging matters. Turning the project into a final report, slide deck, "
        "and submission archive forced the workflow to become more reproducible and more communicable. "
        "That is a valuable outcome on its own because strong cybersecurity work has to be explainable "
        "to other people, not just runnable by its original author."
    )

    add_heading(document, "Conclusion")
    document.add_paragraph(
        "CyberTrace fulfills the project goal by delivering a reproducible Python workflow for "
        "preparing safe cybersecurity data, clustering artifacts, and interpreting the results "
        "through a cybersecurity lens."
    )
    document.save(output_path)


def build_presentation(summary: dict[str, object], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    add_title_slide(
        prs,
        "CyberTrace",
        "Clustering and Analysis of Suspicious Network Traffic\nVictor Andrade and Ben Miller",
    )
    add_bullet_slide(
        prs,
        "Problem Statement",
        [
            "How can suspicious network artifacts be compared without relying only on signatures?",
            "Can unsupervised clustering reveal meaningful groups of benign and malicious behavior?",
            "Can the workflow stay interpretable and useful for cybersecurity investigation?",
        ],
    )
    add_bullet_slide(
        prs,
        "Motivation",
        [
            "Analysts cannot rely only on exact signature matches.",
            "Suspicious traffic often needs comparison by observable behavior.",
            "Clustering helps group similar artifacts for investigation.",
        ],
    )
    add_bullet_slide(
        prs,
        "Contributions",
        [
            "Built a Python repo and CLI for cybersecurity clustering analysis.",
            "Prepared a real NSL-KDD dataset into a reusable feature table.",
            "Generated clustering results, figures, tests, report, slides, and zip packaging.",
        ],
    )
    add_bullet_slide(
        prs,
        "Dataset",
        [
            "Real dataset: small NSL-KDD training CSV sample.",
            f"Records: {summary['record_count']} total, {summary['benign_count']} benign, {summary['malicious_count']} malicious.",
            "Safe benchmark: no live malware and no sensitive production traffic.",
        ],
    )
    add_bullet_slide(
        prs,
        "Pipeline",
        [
            "Raw dataset stored in data/raw/.",
            "prepare-nsl-kdd converts rows into a CyberTrace feature table.",
            "KMeans clusters the standardized numeric feature matrix.",
            "Labels are used after clustering for interpretation only.",
        ],
    )
    add_bullet_slide(
        prs,
        "Methodology",
        [
            "Store raw safe data in data/raw/.",
            "Prepare NSL-KDD into numeric and encoded features.",
            "Standardize the feature matrix and apply KMeans.",
            "Interpret clusters by comparing labels and attack types after fitting.",
        ],
    )
    add_bullet_slide(
        prs,
        "Implementation",
        [
            "Main package: src/cybertrace/",
            "CLI commands: prepare-nsl-kdd and cluster",
            "Outputs: prepared CSV, cluster assignments, notebooks, figures, report",
            "Validation: Ruff checks and Pytest",
        ],
    )
    add_bullet_slide(
        prs,
        "Findings",
        [
            f"Cluster 0: {summary['cluster_rows'][0]['benign']} benign, {summary['cluster_rows'][0]['malicious']} malicious",
            f"Cluster 1: {summary['cluster_rows'][1]['benign']} benign, {summary['cluster_rows'][1]['malicious']} malicious",
            f"Cluster 2: {summary['cluster_rows'][2]['benign']} benign, {summary['cluster_rows'][2]['malicious']} malicious",
        ],
    )
    add_image_slide(
        prs,
        "Results Visualization",
        "PCA projection of the clustered NSL-KDD records.",
        FIGURE_PATH,
    )
    add_bullet_slide(
        prs,
        "Interpretation",
        [
            "Cluster 0 is the most normal-like group.",
            "Cluster 1 is mixed but malicious-heavy, showing overlapping suspicious behavior.",
            "Cluster 2 is fully malicious in this baseline run and dominated by neptune.",
        ],
    )
    add_bullet_slide(
        prs,
        "Implications",
        [
            "Clustering can help prioritize analyst review even without using labels to fit the model.",
            "Mixed clusters are still valuable because they highlight ambiguity worth investigating.",
            "The workflow is reproducible and can be extended to additional safe datasets.",
        ],
    )
    add_bullet_slide(
        prs,
        "Group Members",
        [
            "Victor Andrade: implementation, CLI workflow, tests, packaging, and analysis integration",
            "Ben Miller: cybersecurity framing, project direction, interpretation goals, and final narrative support",
        ],
    )
    add_bullet_slide(
        prs,
        "Future Work",
        [
            "Evaluate newer datasets and richer flow-based artifacts.",
            "Compare KMeans with hierarchical clustering and DBSCAN.",
            "Add deeper per-cluster feature summaries and reporting.",
        ],
    )
    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "CyberTrace shows that clustering can meaningfully support cybersecurity investigation.",
            "The project now delivers a report, presentation, code, data, and a submission zip.",
            "It is a strong foundation for deeper threat-analysis work.",
        ],
    )
    prs.save(output_path)


def add_heading(document: Document, text: str) -> None:
    document.add_heading(text, level=1)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(22)


def add_image_slide(prs: Presentation, title: str, caption: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_box(slide, title)
    slide.shapes.add_picture(str(image_path), PptInches(0.8), PptInches(1.4), width=PptInches(11.7))
    caption_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(6.8), PptInches(11.7), PptInches(0.4))
    paragraph = caption_box.text_frame.paragraphs[0]
    paragraph.text = caption
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = PptPt(18)


def add_title_box(slide, title: str) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0.4),
        PptInches(0.3),
        PptInches(12.5),
        PptInches(0.7),
    )
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = rgb(20, 48, 80)
    box.line.color.rgb = rgb(20, 48, 80)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = PptPt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = rgb(255, 255, 255)


def rgb(red: int, green: int, blue: int):
    from pptx.dml.color import RGBColor

    return RGBColor(red, green, blue)


def assemble_submission_package() -> None:
    if SUBMISSION_DIR.exists():
        shutil.rmtree(SUBMISSION_DIR)
    SUBMISSION_DIR.mkdir(parents=True)

    copy_targets = [
        "README.md",
        "pyproject.toml",
        "Makefile",
        "docs",
        "src",
        "data",
        "reports",
        "tests",
        "notebooks",
        "deliverables",
        "Project Proposal CS544.docx",
    ]
    for target in copy_targets:
        source = ROOT / target
        destination = SUBMISSION_DIR / target
        if source.is_dir():
            shutil.copytree(source, destination, ignore=shutil.ignore_patterns("__pycache__", ".ipynb_checkpoints"))
        else:
            shutil.copy2(source, destination)

    manifest = SUBMISSION_DIR / "SUBMISSION_CONTENTS.txt"
    manifest.write_text(
        "\n".join(
            [
                "CyberTrace submission package",
                "",
                "Included items:",
                "- Final report (deliverables/final_report.md and .docx)",
                "- Presentation slides (deliverables/final_presentation.pptx)",
                "- Source code (src/)",
                "- Data (data/raw and data/processed)",
                "- Results and figures (reports/)",
                "- Notebooks (notebooks/)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


def make_zip_archive() -> None:
    if ZIP_PATH.exists():
        ZIP_PATH.unlink()
    with zipfile.ZipFile(ZIP_PATH, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path in sorted(SUBMISSION_DIR.rglob("*")):
            if path.is_dir():
                continue
            archive.write(path, path.relative_to(SUBMISSION_DIR.parent))


if __name__ == "__main__":
    main()
